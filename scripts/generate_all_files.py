"""
File Generation Script for E2E Testing
Generates all supported file types using appropriate libraries
"""

import os
from pathlib import Path

# Modern Office formats
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

# PDF generation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

# COM automation for legacy Office formats
import win32com.client as win32

def generate_docx():
    """Generate DOCX file using python-docx"""
    print("Generating DOCX file...")
    doc = Document()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('Lorem ipsum dolor sit amet, consectetur adipiscing elit. '
                      'Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. '
                      'Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris.')
    doc.add_paragraph('Duis aute irure dolor in reprehenderit in voluptate velit esse '
                      'cillum dolore eu fugiat nulla pariatur. Excepteur sint occaecat '
                      'cupidatat non proident, sunt in culpa qui officia deserunt mollit anim id est laborum.')
    
    filepath = Path("dev/file_generation/sample.docx")
    doc.save(str(filepath))
    print(f"Created: {filepath}")

def generate_xlsx():
    """Generate XLSX file using openpyxl"""
    print("Generating XLSX file...")
    wb = Workbook()
    ws = wb.active
    ws.title = "Test Data"
    
    # Add headers
    ws['A1'] = 'Name'
    ws['B1'] = 'Description'
    ws['C1'] = 'Value'
    
    # Add sample data
    data = [
        ['Lorem', 'Lorem ipsum dolor sit amet', 123.45],
        ['Ipsum', 'Consectetur adipiscing elit', 678.90],
        ['Dolor', 'Sed do eiusmod tempor incididunt', 234.56],
        ['Amet', 'Ut labore et dolore magna aliqua', 789.01]
    ]
    
    for row_num, row_data in enumerate(data, start=2):
        for col_num, value in enumerate(row_data, start=1):
            ws.cell(row=row_num, column=col_num, value=value)
    
    filepath = Path("dev/file_generation/sample.xlsx")
    wb.save(str(filepath))
    print(f"Created: {filepath}")

def generate_pptx():
    """Generate PPTX file using python-pptx"""
    print("Generating PPTX file...")
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Test Presentation"
    subtitle.text = "Lorem Ipsum Sample Content"
    
    # Content slide
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Sample Content"
    content.text = ("• Lorem ipsum dolor sit amet, consectetur adipiscing elit\n"
                   "• Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua\n"
                   "• Ut enim ad minim veniam, quis nostrud exercitation\n"
                   "• Duis aute irure dolor in reprehenderit in voluptate")
    
    filepath = Path("dev/file_generation/sample.pptx")
    prs.save(str(filepath))
    print(f"Created: {filepath}")

def generate_pdf():
    """Generate PDF file using reportlab"""
    print("Generating PDF file...")
    filepath = Path("dev/file_generation/sample.pdf")
    
    c = canvas.Canvas(str(filepath), pagesize=letter)
    width, height = letter
    
    # Title
    c.setFont("Helvetica-Bold", 16)
    c.drawString(100, height - 100, "Test PDF Document")
    
    # Content
    c.setFont("Helvetica", 12)
    y_position = height - 150
    
    text_lines = [
        "Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do eiusmod",
        "tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim",
        "veniam, quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea",
        "commodo consequat.",
        "",
        "Duis aute irure dolor in reprehenderit in voluptate velit esse cillum",
        "dolore eu fugiat nulla pariatur. Excepteur sint occaecat cupidatat non",
        "proident, sunt in culpa qui officia deserunt mollit anim id est laborum.",
        "",
        "Sed ut perspiciatis unde omnis iste natus error sit voluptatem accusantium",
        "doloremque laudantium, totam rem aperiam, eaque ipsa quae ab illo inventore",
        "veritatis et quasi architecto beatae vitae dicta sunt explicabo."
    ]
    
    for line in text_lines:
        c.drawString(100, y_position, line)
        y_position -= 20
    
    c.save()
    print(f"Created: {filepath}")

def generate_doc_com():
    """Generate DOC file using COM automation with Microsoft Word"""
    print("Generating DOC file using COM automation...")
    
    try:
        # Start Word application
        word = win32.Dispatch("Word.Application")
        word.Visible = False  # Don't show Word window
        
        # Create new document
        doc = word.Documents.Add()
        
        # Add title
        doc.Range().Font.Size = 16
        doc.Range().Font.Bold = True
        doc.Range().InsertAfter("Test Legacy DOC Document\n\n")
        
        # Add content
        doc.Range().Font.Size = 12
        doc.Range().Font.Bold = False
        content = """Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea commodo consequat.

Duis aute irure dolor in reprehenderit in voluptate velit esse cillum dolore eu fugiat nulla pariatur. Excepteur sint occaecat cupidatat non proident, sunt in culpa qui officia deserunt mollit anim id est laborum.

Sed ut perspiciatis unde omnis iste natus error sit voluptatem accusantium doloremque laudantium, totam rem aperiam, eaque ipsa quae ab illo inventore veritatis et quasi architecto beatae vitae dicta sunt explicabo."""
        
        doc.Range().InsertAfter(content)
        
        # Save as DOC format (Word 97-2003)
        filepath = os.path.abspath("dev/file_generation/sample.doc")
        doc.SaveAs(filepath, FileFormat=0)  # 0 = Word 97-2003 Document
        
        # Close document and quit Word
        doc.Close()
        word.Quit()
        
        print(f"Created: {filepath}")
        
    except Exception as e:
        print(f"Error creating DOC file: {e}")
        print("Make sure Microsoft Word is installed and accessible.")

def generate_xls_com():
    """Generate XLS file using COM automation with Microsoft Excel"""
    print("Generating XLS file using COM automation...")
    
    try:
        # Start Excel application
        excel = win32.Dispatch("Excel.Application")
        excel.Visible = False  # Don't show Excel window
        
        # Create new workbook
        wb = excel.Workbooks.Add()
        ws = wb.ActiveSheet
        ws.Name = "Test Data"
        
        # Add headers
        ws.Cells(1, 1).Value = "Name"
        ws.Cells(1, 2).Value = "Description" 
        ws.Cells(1, 3).Value = "Value"
        
        # Format headers
        header_range = ws.Range("A1:C1")
        header_range.Font.Bold = True
        header_range.Interior.ColorIndex = 15  # Light gray
        
        # Add sample data
        data = [
            ['Lorem', 'Lorem ipsum dolor sit amet', 123.45],
            ['Ipsum', 'Consectetur adipiscing elit', 678.90],
            ['Dolor', 'Sed do eiusmod tempor incididunt', 234.56],
            ['Amet', 'Ut labore et dolore magna aliqua', 789.01]
        ]
        
        for row_num, row_data in enumerate(data, start=2):
            for col_num, value in enumerate(row_data, start=1):
                ws.Cells(row_num, col_num).Value = value
        
        # Auto-fit columns
        ws.Columns.AutoFit()
        
        # Save as XLS format (Excel 97-2003)
        filepath = os.path.abspath("dev/file_generation/sample.xls")
        wb.SaveAs(filepath, FileFormat=56)  # 56 = Excel 97-2003 Workbook
        
        # Close workbook and quit Excel
        wb.Close()
        excel.Quit()
        
        print(f"Created: {filepath}")
        
    except Exception as e:
        print(f"Error creating XLS file: {e}")
        print("Make sure Microsoft Excel is installed and accessible.")

def generate_ppt_com():
    """Generate PPT file using COM automation with Microsoft PowerPoint"""
    print("Generating PPT file using COM automation...")
    
    try:
        # Start PowerPoint application
        powerpoint = win32.Dispatch("PowerPoint.Application")
        # PowerPoint doesn't support Visible = False, so we leave it as default
        
        # Create new presentation
        prs = powerpoint.Presentations.Add()
        
        # Add title slide
        slide1 = prs.Slides.Add(1, 1)  # ppLayoutTitle
        slide1.Shapes.Title.TextFrame.TextRange.Text = "Test Legacy PPT Presentation"
        slide1.Shapes(2).TextFrame.TextRange.Text = "Lorem Ipsum Sample Content"
        
        # Add content slide
        slide2 = prs.Slides.Add(2, 2)  # ppLayoutText
        slide2.Shapes.Title.TextFrame.TextRange.Text = "Sample Content"
        slide2.Shapes(2).TextFrame.TextRange.Text = """• Lorem ipsum dolor sit amet, consectetur adipiscing elit
• Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua
• Ut enim ad minim veniam, quis nostrud exercitation
• Duis aute irure dolor in reprehenderit in voluptate"""
        
        # Save as PPT format (PowerPoint 97-2003)
        filepath = os.path.abspath("dev/file_generation/sample.ppt")
        prs.SaveAs(filepath, 1)  # 1 = PowerPoint 97-2003 Presentation
        
        # Close presentation and quit PowerPoint
        prs.Close()
        powerpoint.Quit()
        
        print(f"Created: {filepath}")
        
    except Exception as e:
        print(f"Error creating PPT file: {e}")
        print("Make sure Microsoft PowerPoint is installed and accessible.")

def main():
    """Generate all file types"""
    print("Starting file generation for E2E testing...\n")
    
    # Create output directory if it doesn't exist
    output_dir = Path("dev/file_generation")
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Generate modern Office formats
    print("=== Modern Office Formats ===")
    generate_docx()
    generate_xlsx()
    generate_pptx()
    
    # Generate PDF
    print("\n=== PDF Format ===")
    generate_pdf()
    
    # Generate legacy Office formats using COM
    print("\n=== Legacy Office Formats (COM Automation) ===")
    generate_doc_com()
    generate_xls_com()
    generate_ppt_com()
    
    print("\n=== File Generation Complete ===")
    print("All files created in dev/file_generation/")
    
    # List created files
    print("\nGenerated files:")
    for file_path in output_dir.glob("*"):
        if file_path.is_file():
            size_kb = file_path.stat().st_size / 1024
            print(f"  {file_path.name} ({size_kb:.1f} KB)")

if __name__ == "__main__":
    main()